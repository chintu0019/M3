"""
M3 Content Extractors -- extract text from URLs, PDFs, and other content types.

All heavy/sync operations run via asyncio.to_thread.
"""

import asyncio
import logging

import httpx

logger = logging.getLogger("m3.extractors")


async def extract_url(url: str) -> str:
    """Fetch a URL and extract readable text content."""
    async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
        response = await client.get(url)
        response.raise_for_status()

    import trafilatura

    text = await asyncio.to_thread(trafilatura.extract, response.text)
    return text or ""


async def extract_pdf(pdf_bytes: bytes) -> str:
    """Extract text from PDF bytes."""

    def _extract():
        import io

        import pdfplumber

        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            return "\n\n".join(page.extract_text() or "" for page in pdf.pages)

    return await asyncio.to_thread(_extract)


async def extract_docx(docx_bytes: bytes) -> str:
    """Extract text from a .docx file."""

    def _extract():
        import io

        from docx import Document

        doc = Document(io.BytesIO(docx_bytes))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)

    return await asyncio.to_thread(_extract)


async def extract_xlsx(xlsx_bytes: bytes) -> str:
    """Extract text from an .xlsx spreadsheet."""

    def _extract():
        import io

        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(xlsx_bytes), data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"# Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
            parts.append("")
        wb.close()
        return "\n".join(parts)

    return await asyncio.to_thread(_extract)


async def extract_pptx(pptx_bytes: bytes) -> str:
    """Extract text from a .pptx presentation."""

    def _extract():
        import io

        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_parts.append(f"Notes: {slide.notes_slide.notes_text_frame.text}")
            parts.append("\n\n".join(slide_parts))
        return "\n\n".join(parts)

    return await asyncio.to_thread(_extract)


async def extract_epub(epub_bytes: bytes) -> str:
    """Extract text from an .epub ebook."""

    def _extract():
        import io

        import trafilatura
        from ebooklib import ITEM_DOCUMENT, epub

        book = epub.read_epub(io.BytesIO(epub_bytes))
        parts = []
        for item in book.get_items_of_type(ITEM_DOCUMENT):
            html = item.get_content().decode("utf-8", errors="replace")
            text = trafilatura.extract(html) or ""
            if text.strip():
                parts.append(text)
        return "\n\n".join(parts)

    return await asyncio.to_thread(_extract)


async def extract_html(html_bytes: bytes) -> str:
    """Extract readable text from HTML content."""

    def _extract():
        import trafilatura

        html = html_bytes.decode("utf-8", errors="replace")
        return trafilatura.extract(html) or html

    return await asyncio.to_thread(_extract)


TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv",
    ".json", ".jsonl", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".xml", ".srt", ".vtt",
    # Code
    ".py", ".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs",
    ".rb", ".go", ".rs", ".java", ".kt", ".swift",
    ".c", ".cpp", ".cc", ".h", ".hpp", ".cs",
    ".php", ".pl", ".r", ".scala", ".clj", ".ex", ".exs", ".elm",
    ".sh", ".bash", ".zsh", ".fish", ".ps1",
    ".sql", ".graphql", ".gql",
    ".css", ".scss", ".sass", ".less",
    ".dockerfile",
    ".env", ".envrc",
}


async def extract_by_filename(file_bytes: bytes, filename: str) -> str:
    """Dispatch extraction by filename extension for generic 'file' content type."""
    name = filename.lower()

    if name.endswith(".pdf"):
        return await extract_pdf(file_bytes)
    if name.endswith(".docx"):
        return await extract_docx(file_bytes)
    if name.endswith(".xlsx"):
        return await extract_xlsx(file_bytes)
    if name.endswith(".pptx"):
        return await extract_pptx(file_bytes)
    if name.endswith(".epub"):
        return await extract_epub(file_bytes)
    if name.endswith((".html", ".htm")):
        return await extract_html(file_bytes)

    # Text-like files
    for ext in TEXT_EXTENSIONS:
        if name.endswith(ext):
            try:
                return file_bytes.decode("utf-8", errors="replace")
            except Exception:
                return ""

    # Last-ditch attempt: try to decode as text if content looks textual
    try:
        decoded = file_bytes.decode("utf-8")
        # Heuristic: if it mostly decodes cleanly and has reasonable character density, treat as text
        if len(decoded) > 0 and sum(1 for c in decoded[:1000] if c.isprintable() or c in "\n\r\t") / min(len(decoded), 1000) > 0.9:
            return decoded
    except UnicodeDecodeError:
        pass

    return ""
